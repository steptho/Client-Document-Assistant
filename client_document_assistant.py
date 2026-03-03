# Client Document Assistant

import os
import tempfile
import base64
import json
import streamlit as st
import pandas as pd
from fpdf import FPDF
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openai import OpenAI
from striprtf.striprtf import rtf_to_text
import io

# ---------------------------------------------------
# CONFIGURATION & PERSISTENCE
# ---------------------------------------------------
st.set_page_config(page_title="Document Data Assistant", layout="wide")
HISTORY_FILE = "analysis_history.json"

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r") as f:
                return json.load(f)
        except:
            return []
    return []

def save_to_history(new_item):
    history = load_history()
    history.append(new_item)
    with open(HISTORY_FILE, "w") as f:
        json.dump(history, f, indent=4)

# ---------------------------------------------------
# SESSION STATE INITIALIZATION
# ---------------------------------------------------
if "state" not in st.session_state:
    st.session_state.state = {
        "logged_in": False,
        "analysis_result": "",
        "analysis_history": load_history(),
        "preview_text": ""
    }

# ---------------------------------------------------
# OPENAI CLIENT
# ---------------------------------------------------
try:
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
except:
    st.error("Please set your OPENAI_API_KEY in Streamlit Secrets.")

# ---------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------

def extract_text(uploaded_file):
    """Reads content from various file types including AI Vision for images."""
    # Use splitext to get the extension correctly
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    try:
        uploaded_file.seek(0)
        
        # --- 1. PDF ---
        if ext == ".pdf":
            reader = PdfReader(uploaded_file)
            text = "\n".join([p.extract_text() or "" for p in reader.pages])
        
        # --- 2. WORD (.DOCX) ---
        elif ext == ".docx":
            doc = Document(uploaded_file)
            text = "\n".join([p.text for p in doc.paragraphs])
            
        # --- 3. LEGACY WORD (.DOC) ---
        elif ext == ".doc":
            text = (f"⚠️ **Legacy Format Detected:** '{uploaded_file.name}' is an older .doc file.\n\n"
                    f"**Action Required:** Please open this file in Word, select 'Save As', "
                    f"and choose **Word Document (.docx)**. Then, upload the new version here.")
        # --- 4. POWERPOINT ---
        elif ext == ".pptx":
            prs = Presentation(uploaded_file)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        
        # --- 5. SPREADSHEETS (CSV, XLSX, XLSM, XLS) ---
        elif ext in [".xlsx", ".xls", ".xlsm", ".csv"]:
            if ext == ".csv":
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            
            # 1. Drop columns that are completely empty
            df = df.dropna(axis=1, how='all')

            # 2. Truncate any cell that has more than 500 characters 
            # (Prevents one giant cell from breaking the token limit)
            df = df.map(lambda x: str(x)[:500] if isinstance(x, str) else x)

            # 3. Final safety check on total size
            full_text = df.to_string(index=False)
            if len(full_text) > 100000: # Approx 25k tokens
                text = (f"Note: Data is very dense. Showing first 100 rows only.\n"
                        f"{df.head(100).to_string(index=False)}")
            else:
                text = f"Spreadsheet Data ({uploaded_file.name}):\n{full_text}"

        # --- 6. RTF ---
        elif ext == ".rtf":
            content = uploaded_file.read().decode('utf-8', errors='ignore')
            text = rtf_to_text(content)

        # --- 7. IMAGES ---
        elif ext in [".jpg", ".jpeg", ".png"]:
            file_bytes = uploaded_file.read()
            base64_image = base64.b64encode(file_bytes).decode("utf-8")
            
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a professional analyst. Describe the scene or extract text precisely."},
                    {"role": "user", "content": [
                        {"type": "text", "text": "Describe this image in detail or transcribe its content."},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}", "detail": "high"}}
                    ]}
                ]
            )
            text = response.choices[0].message.content
            
    except Exception as e:
        st.error(f"Error reading file {uploaded_file.name}: {e}")
    return text

def generate_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    clean_text = text.replace('’', "'").replace('“', '"').replace('”', '"').replace('—', '-')
    clean_text = clean_text.encode('latin-1', 'ignore').decode('latin-1')
    pdf.multi_cell(190, 10, clean_text, align='L')
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        pdf.output(tmp.name)
        return tmp.name

def handle_login():
    if st.session_state.get("temp_user_input", "").strip():
        st.session_state.state["logged_in"] = True

def clear_old_file_state():
    st.session_state.state["preview_text"] = ""
    st.session_state.state["analysis_result"] = ""

# ---------------------------------------------------
# MAIN UI LOGIC
# ---------------------------------------------------

if not st.session_state.state["logged_in"]:
    st.title("🔐 Secure Login")
    st.text_input("Enter Company Name", key="temp_user_input", on_change=handle_login)
    if st.button("Enter"):
        handle_login()
        st.rerun()

else:
    st.session_state.state["analysis_history"] = load_history()
    
    with st.sidebar:
        st.header("🏢 Session")
        # ... logout button ...
        
        st.divider()
        st.subheader("📜 History")
        
        # ALWAYS pull the latest from the session state here
        history_list = st.session_state.state["analysis_history"]
        
        # ... search and summarize logic ...
        search_query = st.text_input("🔍 Search History", "").lower()
        history_list = st.session_state.state["analysis_history"]
        
        if history_list:
            if st.button("🤖 Summarize All History"):
                with st.spinner("Analyzing history..."):
                    all_text = "\n\n".join([f"Document: {h['title']}\nContent: {h['analysis']}" for h in history_list])
                    sum_res = client.chat.completions.create(
                        model="gpt-4o-mini",
                        messages=[{"role": "user", "content": f"Provide a categorized summary and highlights of these records:\n{all_text}"}]
                    )
                    st.session_state.state["analysis_result"] = f"### 📊 GLOBAL HISTORY SUMMARY\n\n{sum_res.choices[0].message.content}"
                    st.rerun()

            st.divider()
            filtered_history = [item for item in history_list if search_query in item['title'].lower() or search_query in item['analysis'].lower()]
            for idx, item in enumerate(filtered_history):
                if st.button(f"📄 {item['title']}", key=f"hist_{idx}"):
                    st.session_state.state["analysis_result"] = item["analysis"]
                    st.session_state.state["preview_text"] = ""
                    st.rerun()

    st.title("📂 Document Data Assistant")

    # --- THE UPDATED BOUNCER (Uploader) ---
    uploaded_file = st.file_uploader(
        "Upload PDF, DOCX, DOC, PPTX, XLSX, XLSM, CSV, or Images", 
        type=["pdf", "docx", "doc", "pptx", "xlsx", "xlsm", "csv", "jpg", "jpeg", "png", "rtf"],
        on_change=clear_old_file_state
    )

    if uploaded_file:
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        is_image = file_ext in [".jpg", ".jpeg", ".png"]

        if not st.session_state.state["preview_text"]:
            with st.spinner("🔍 AI is reading the file..."):
                st.session_state.state["preview_text"] = extract_text(uploaded_file)

        st.subheader("📄 Document Preview")
        if is_image:
            st.image(uploaded_file, width=500)
            with st.expander("See AI-Extracted Raw Data"):
                st.write(st.session_state.state["preview_text"])
        else:
            st.text_area("Content", st.session_state.state["preview_text"], height=200)

        if st.button("🚀 Run Full Analysis"):
            raw_content = st.session_state.state["preview_text"]
            
            # STOP if there's no content or if it's just the .doc warning
            if not raw_content or raw_content.startswith("⚠️ Note:"):
                st.warning("No readable content found to analyze.")
            else:
                with st.spinner("📊 AI is generating your report..."):
                    try:
                        # 1. Generate Title
                        t_res = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[{"role": "user", "content": f"4-word title for: {raw_content[:500]}"}]
                        )
                        # We define 'title' here
                        title = t_res.choices[0].message.content.strip().replace('"', '')

                        # 2. Generate Analysis
                        a_res = client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[{"role": "user", "content": f"Summarize key insights and data from this text:\n{raw_content}"}]
                        )
                        # We define 'analysis' here
                        analysis = a_res.choices[0].message.content

                        # 3. SAVE ONLY IF BOTH EXIST (Inside the try block)
                        save_to_history({"title": title, "analysis": analysis})
                        
                        # 4. Update memory and UI
                        st.session_state.state["analysis_history"] = load_history()
                        st.session_state.state["analysis_result"] = analysis
                        st.rerun()

                    except Exception as e:
                        st.error(f"Something went wrong during AI analysis: {e}")

    if st.session_state.state["analysis_result"]:
        st.divider()
        st.subheader("💡 Analysis Findings")
        st.markdown(st.session_state.state["analysis_result"])
        pdf_path = generate_pdf(st.session_state.state["analysis_result"])
        with open(pdf_path, "rb") as f:
            st.download_button(label="📥 Download PDF", data=f, file_name="Analysis_Report.pdf", mime="application/pdf")
